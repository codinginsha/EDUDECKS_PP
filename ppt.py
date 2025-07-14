import google.generativeai as genai
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import os
import random
import re
import string
from googletrans import Translator
from googleapiclient.discovery import build
import datetime

# === API KEYS ===
GENAI_API_KEY = os.environ.get("GENAI_API_KEY")
UNSPLASH_ACCESS_KEY_1 = os.environ.get("UNSPLASH_ACCESS_KEY_1")
#UNSPLASH_ACCESS_KEY_2 = "VdTXCqfdaXOoxazESWkqmtIszRfFFWZJEBO61ZaS_80"
YOUTUBE_API_KEY = os.environ.get("YOUTUBE_API_KEY")

if not GENAI_API_KEY:
    print("⚠️ Warning: GENAI_API_KEY environment variable not set!")
if not UNSPLASH_ACCESS_KEY_1:
    print("⚠️ Warning: UNSPLASH_ACCESS_KEY_1 environment variable not set!")
if not YOUTUBE_API_KEY:
    print("⚠️ Warning: YOUTUBE_API_KEY environment variable not set!")

# === Configuration ===
genai.configure(api_key=GENAI_API_KEY)
translator = Translator()

BACKGROUND_COLORS = [
    "#FAFAFA",  # Soft White
    "#F5F5F5",  # Off-White
    "#EAF4FB",  # Light Pastel Blue
    "#DDEAF6",  # Gentle Sky Blue
    "#F9E6E6",  # Dusty Rose
    "#FCEEEF",  # Blush Pink
    "#FFF8E7",  # Pale Beige / Cream
    "#FDF6EC",  # Warm Cream
    "#E8F0E3",  # Muted Sage Green
    "#DDEBDD",  # Calm Green Tint
    "#F0F0F0",  # Very Light Gray
    "#EBEBEB",  # Ultra Light Gray
    "#F4F1FA",  # Lavender Mist
    "#EDEAF6",  # Soft Lilac
    "#E6FFFA",  # Icy Mint
    "#D2F4EA",  # Minty Aqua
    "#F8EFD4",  # Champagne Beige
]

# Title fonts - bold, decorative fonts
TITLE_FONT_MAP = {
    "urdu": [
        "Jameel Noori Nastaleeq", "Alvi Nastaleeq", "Mehr Nastaliq", "Fajer Noori Nastalique", "Nafees Nastaleeq", "Awami Nastaliq"
    ],
    "marathi": [
        "Mangal", "Kokila", "Aparajita", "Sanskrit Text", "Utsaah", "Shruti", "Sahadeva", "Nirmala UI", "Devanagari MT", "Devanagari Sangam MN"
    ],
    "hindi": [
        "Mangal", "Samarkan", "Kokila", "Aparajita", "Sanskrit Text", "Utsaah", "Shruti", "Sahadeva", "Nirmala UI", "Devanagari MT", "Devanagari Sangam MN"
    ],
    "english": [
        "Montserrat", "Poppins", "Lobster", "Oswald", "Raleway", "Arial Black", "Century Gothic", "Franklin Gothic Medium", "Futura", "Bebas Neue", "Brush Script MT", "Copperplate Gothic Bold"
    ]
}

# Content fonts - clean, readable fonts
CONTENT_FONT_MAP = {
    "urdu": [
        "Jameel Noori Nastaleeq", "Alvi Nastaleeq", "Mehr Nastaliq", "Fajer Noori Nastalique", "Nafees Nastaleeq", "Awami Nastaliq"
    ],
    "marathi": [
        "Mangal", "Kokila", "Aparajita", "Sanskrit Text", "Utsaah", "Shruti", "Sahadeva", "Nirmala UI", "Devanagari MT", "Devanagari Sangam MN"
    ],
    "hindi": [
        "Mangal", "Samarkan", "Kokila", "Aparajita", "Sanskrit Text", "Utsaah", "Shruti", "Sahadeva", "Nirmala UI", "Devanagari MT", "Devanagari Sangam MN"
    ],
    "english": [
        "Poppins", "Montserrat", "Raleway", "Lato", "Nunito", "Quicksand", "Segoe UI", "Calibri", "Arial", "Verdana", "Tahoma", "Candara", "Lucida Sans Unicode", "Georgia", "Garamond", "Trebuchet MS", "Gill Sans", "Perpetua"
    ]
}

# Add after CONTENT_FONT_MAP
CONTENT_FONT_SCALING = {
    # Sans-serif modern fonts (tend to look larger)
    "Poppins": 0.92,
    "Montserrat": 0.92,
    "Raleway": 0.94,
    "Lato": 0.96,
    "Nunito": 0.96,
    "Quicksand": 0.95,
    # System/standard sans-serif
    "Segoe UI": 1.0,
    "Calibri": 1.0,
    "Arial": 1.0,
    "Verdana": 0.98,
    "Tahoma": 1.0,
    "Candara": 1.0,
    "Lucida Sans Unicode": 1.02,
    # Serif fonts (tend to look smaller)
    "Georgia": 1.08,
    "Garamond": 1.10,
    "Trebuchet MS": 1.0,
    "Gill Sans": 1.0,
    "Perpetua": 1.10
}

# Fallback list for random selection if language not found
FONT_LIST = [
    "Montserrat", "Poppins", "Lobster", "Oswald", "Raleway", "Comic Sans MS",
    "Georgia", "Garamond", "Trebuchet MS", "Segoe UI", "Calibri",
    "Arial Black", "Century Gothic", "Franklin Gothic Medium", "Futura", "Bebas Neue", "Brush Script MT", "Palatino Linotype", "Copperplate Gothic Bold",
    "Verdana", "Tahoma", "Candara", "Lucida Sans Unicode", "Rockwell", "Baskerville", "Gill Sans", "Perpetua", "Courier New"
]

BACKGROUND_IMAGE_PATHS = [
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Pink and Green Doodle Hand drawn Science Project Presentation.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Purple Blue and Pink Playful Illustrative Sermon Church Presentation.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Teal  White Playful Creative Project Presentation (1).jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Teal  White Playful Creative Project Presentation.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\WhatsApp Image 2025-07-14 at 00.51.19_e062a5a1.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\2.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\10.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\44.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Beige Dark Grey Vintage Victorian Project History Presentation.png",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Brown Illustrative Playful Ancient Egypt History Presentation.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Brown Illustrative Playful Ancient Egypt History Presentation.png",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\2.png",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0020.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0021.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0022.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0023.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0011.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0012.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0013.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0014.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0015.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0016.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0007.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0009.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0031.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0024.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0025.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0026.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0027.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0028.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0029.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0030.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0017.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0018.jpg",
    r"C:\Users\momin\Downloads\ppt templates\IMG-20250714-WA0019.jpg",
    r"C:\Users\momin\Downloads\ppt templates\Screenshot 2025-07-14 161839.png",
    r"C:\Users\momin\Downloads\ppt templates\Screenshot 2025-07-14 161806.png",
    r"C:\Users\momin\Downloads\ppt templates\Screenshot 2025-07-14 161936.png",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0017.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0018.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0019.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0020.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0021.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0022.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0023.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0024.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0025.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0027.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0028.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0029.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0030.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0031.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0032.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0033.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0034.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0035.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0036.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Screenshot 2025-07-14 161806.png",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Screenshot 2025-07-14 161839.png",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Screenshot 2025-07-14 161936.png",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0007.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0009.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0011.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0012.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0013.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0014.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0015.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0016.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0036.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0030.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0032.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0033.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0034.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0035.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Pink and Blue Gradient Annual Report Presentation.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Pink Purple Modern Minimalist Aesthetic Project Presentation.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\_ੈ✩‧₊˚ ʙʀɪɢʜᴛ ᴘɪɴᴋ ᴘᴏᴡᴇʀᴘᴏɪɴᴛ ᴛᴇᴍᴘʟᴀᴛᴇ ༊_·˚.jpeg.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\download (12).jpeg.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\download (13).jpeg.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Fundos para apresentações.jpeg.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Gradient Colorful Minimalist Coming  Soon Banner.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Gray and White Simple Thesis Defense Presentation.jpg",
    r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Minimalist & Modern Fashion Portfolio Presentation 1.jpg",
]


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def is_dark_color(rgb_color):
    r, g, b = rgb_color
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    return brightness < 128


def sanitize_filename(filename):
    valid_chars = f"-_.() {string.ascii_letters}{string.digits}"
    return ''.join(c for c in filename if c in valid_chars)


def translate_to_english(text, lang='auto'):
    try:
        translated = translator.translate(text, src=lang, dest='en')
        return translated.text
    except Exception as e:
        print(f"\u274c Translation error: {e}")
        return text


def fetch_unsplash_image(subject, query, language='English'):
    search_term = f"{subject} {query}"
    if language.lower() != 'english':
        search_term = translate_to_english(search_term, lang=language.lower())
    search_term = search_term.strip().replace(" ", "+")

    url = f"https://api.unsplash.com/photos/random?query={search_term}&orientation=landscape&client_id={UNSPLASH_ACCESS_KEY_1}"
    try:
        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            image_url = data['urls']['regular']
            print(f"\u2705 Found image for: {search_term}")
            img_data = requests.get(image_url).content
            return BytesIO(img_data)
        else:
            print(f"\u26a0\ufe0f Unsplash error: {response.status_code} - {response.text}")
    except Exception as e:
        print(f"\u274c Image fetch error: {e}")
    return None


def search_youtube_video(query):
    # Return a YouTube search results URL for the query
    search_query = query.replace(' ', '+')
    return f"https://www.youtube.com/results?search_query={search_query}"


def split_into_bullets(text):
    lines = text.split('\n')
    bullets = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("\u2022") or line.startswith("-"):
            clean = re.sub(r"^[-\u2022\s]+", "", line)
            bullets.append(f"\u2022 {clean}")
        else:
            parts = re.split(r'(?<=[.!?])\s+', line)
            for part in parts:
                part = part.strip()
                if part:
                    bullets.append(f"\u2022 {part}")
    return bullets


def clean_title(title, language=None):
    import re
    # Remove asterisks, special characters, and formatting artifacts for both display and filenames
    title = re.sub(r'[\\/*?<>|:"\n]', '', title)
    title = title.replace('<u>', '').replace('</u>', '').strip()
    # Remove stray 'u' or 'U' at the start/end (for Urdu)
    if language and language.lower() == 'urdu':
        # Remove leading/trailing 'u' or 'U' (with or without spaces)
        title = re.sub(r'^[uU]\s*', '', title)
        title = re.sub(r'\s*[uU]$', '', title)
        # Remove any leading/trailing English letters or symbols
        title = re.sub(r'^[A-Za-z\s\-\_\,\.\!\?\(\)\[\]]+', '', title)
        title = re.sub(r'[A-Za-z\s\-\_\,\.\!\?\(\)\[\]]+$', '', title)
        # Remove any remaining English words
        title = re.sub(r'[A-Za-z]+', '', title)
        title = title.strip()
    else:
        # Remove stray 'u' at the start if present (from parsing)
        if title.lower().startswith('u') and len(title) > 1 and title[1].isupper():
            title = title[1:]
        # Remove stray 'u' at the end if present (for RTL/Urdu)
        if title.endswith('u') and len(title) > 2 and not title[-2].isspace():
            title = title[:-1]
        # Remove stray 'u' at the end if it is separated by a space (for RTL/Urdu)
        if title.endswith(' u'):
            title = title[:-2]
    return title.strip()


def generate_ppt_content(class_level, subject, topic, language='English', num_slides=5):
    model = genai.GenerativeModel('gemini-2.0-flash')
    is_language_subject = subject.strip().lower() in ["english", "hindi", "urdu", "marathi"]
    is_poem_or_lesson = any(x in topic.strip().lower() for x in ["poem", "lesson"])
    prompt = f"""Create a {num_slides}-slide PowerPoint presentation for a {class_level}th grade {subject} class on the topic \"{topic}\".
Use only formal, culturally accurate, and age-appropriate {language} language.
Do not mix English with {language}. No Roman script or SMS-style writing.
Each slide should include:
- A meaningful, localized title (bold and underlined using *<u>Title</u>* format).
- 5-6 clear bullet points using (\u2022 or -), with examples or facts.
- Do not write paragraphs, just concise bullet points.
- Avoid any unrelated historical figures or general knowledge.

If the subject or topic is related to math, mathematics, numericals, or formulas, include relevant formulas and worked examples on the slides. Use proper Unicode math formatting (e.g., x², fractions, π, √, etc.) and avoid ASCII-only equations. Show at least one worked example for each formula if possible.

"""
    if is_language_subject and is_poem_or_lesson:
        prompt += f"\nAfter the content slides, add one slide titled 'Difficult Words and Meanings' (in {language}) with 5-8 hard words from the poem or lesson and their meanings as bullet points. Format: \u2022 word: meaning."
    prompt += f"\nAfter the {{num_slides}} content slides, add one practice slide with 5 questions."
    prompt += f"\n\nRespond exactly in this format:\n\nSLIDE 1: [Title]\n[\u2022 Bullet or - Bullet or sentence]\n\nSLIDE 2: [Title]\n[\u2022 Bullet]\n[...]\n\nSLIDE {{num_slides + 1}}: Practice Questions\n[\u2022 Question 1]\n[\u2022 Question 2]\n[\u2022 Question 3]\n[\u2022 Question 4]\n[\u2022 Question 5]"
    try:
        response = model.generate_content(prompt)
        return format_ppt_output(response.text)
    except Exception as e:
        return f"\u274c Error generating content: {str(e)}"


def format_ppt_output(raw_text):
    slides = []
    current_slide = {}
    for line in raw_text.split('\n'):
        if line.strip().startswith('SLIDE'):
            if current_slide:
                slides.append(current_slide)
            parts = line.split(':', 1)
            current_slide = {'title': clean_title(parts[1].strip()), 'content': []}
        elif line.strip():
            current_slide['content'].append(line.strip())
    if current_slide:
        slides.append(current_slide)
    return slides


def set_slide_background_image(slide, prs, image_path):
    """Set a background image for a slide."""
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    # Remove any existing background fill
    slide.background.fill.solid()
    # Add the image as a picture shape covering the whole slide
    pic = slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)
    # Send the image to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

def get_contrasting_font_color(image_path):
    """Determine the best font color based on background image brightness."""
    try:
        from PIL import Image
        import numpy as np
        
        # Open and analyze the image
        with Image.open(image_path) as img:
            # Convert to RGB if needed
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Resize for faster processing
            img = img.resize((100, 100))
            
            # Convert to numpy array and calculate average brightness
            img_array = np.array(img)
            avg_brightness = np.mean(img_array)
            
            # Determine font color based on background brightness
            if avg_brightness > 128:  # Light background
                return RGBColor(0, 0, 0)  # Black text
            else:  # Dark background
                return RGBColor(255, 255, 255)  # White text
                
    except Exception as e:
        print(f"⚠️ Could not analyze image brightness: {e}")
        return RGBColor(0, 0, 0)  # Default to black


def create_powerpoint(slides, topic, class_level, subject, language="English"):
    prs = Presentation()
    language_key = language.lower().strip()
    rtl = True if language_key == 'urdu' else False

    yt_query = f"{subject} {topic} {language} grade {class_level}"
    yt_url = search_youtube_video(yt_query)

    # Randomly select a background image for this presentation
    # Exclude certain images for Urdu language
    urdu_excluded_images = [
        r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Pink and Green Doodle Hand drawn Science Project Presentation.jpg",
        r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Green and Orange Playful Kids Birthday Presentation.jpg",
        r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\Teal  White Playful Creative Project Presentation.jpg",
        r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\44.jpg"
    ]
    
    english_excluded_images = [
        r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\1.png",
        r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\2.png",
        r"C:\Users\momin\OneDrive\Documents\EduDECKS_PPT\ppt templates\IMG-20250714-WA0025.jpg"
    ]
    if language_key == 'urdu':
        # Filter out excluded images for Urdu
        available_images = [img for img in BACKGROUND_IMAGE_PATHS if img not in urdu_excluded_images]
        background_image_path = random.choice(available_images)
    elif language_key == 'english':
        # Filter out excluded images for English
        available_images = [img for img in BACKGROUND_IMAGE_PATHS if img not in english_excluded_images]
        background_image_path = random.choice(available_images)
    else:
        background_image_path = random.choice(BACKGROUND_IMAGE_PATHS)
    # Pick fonts for this presentation
    title_font_choices = TITLE_FONT_MAP.get(language_key, FONT_LIST)
    content_font_choices = CONTENT_FONT_MAP.get(language_key, FONT_LIST)
    title_font = random.choice(title_font_choices)
    content_font = random.choice(content_font_choices)
    # Get contrasting font color based on background image
    font_color = get_contrasting_font_color(background_image_path)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background_image(title_slide, prs, background_image_path)

    # --- Title Slide Centered Title and Subtitle (No Overlap) ---
    first_title = clean_title(slides[0]['title'], language)
    slide_width = float(prs.slide_width) / 914400 if prs.slide_width else 10.0
    slide_height = float(prs.slide_height) / 914400 if prs.slide_height else 7.5

    # One big vertical box in the center
    box_height = 3.5  # inches, adjust as needed
    box_top = (slide_height - box_height) / 2.0
    title_box = title_slide.shapes.add_textbox(
        Inches(1.0),
        Inches(box_top),
        Inches(slide_width - 2.0),
        Inches(box_height)
    )
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Title paragraph
    p = tf.add_paragraph()
    p.text = first_title
    p.font.bold = True
    p.font.name = title_font
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.CENTER
    title_len = len(first_title)
    if title_len > 120:
        p.font.size = Pt(36)
    elif title_len > 80:
        p.font.size = Pt(44)
    elif title_len > 50:
        p.font.size = Pt(50)
    else:
        p.font.size = Pt(54)

    # Subtitle paragraph (if present)
    if slides[0]['content']:
        subtitle = slides[0]['content'][0]
        sp = tf.add_paragraph()
        sp.text = subtitle[:250]
        sp.font.size = Pt(32) if len(subtitle) < 80 else Pt(24)
        sp.font.color.rgb = font_color
        sp.font.name = content_font
        sp.alignment = PP_ALIGN.CENTER
        sp.space_before = Pt(24)  # extra space between title and subtitle

    # Place YouTube link at the bottom of the title slide
    if yt_url:
        yt_box = title_slide.shapes.add_textbox(
            Inches(1.0),
            Inches(slide_height - 0.7),
            Inches(slide_width - 2.0),
            Inches(0.5)
        )
        yt_tf = yt_box.text_frame
        yt_tf.clear()
        p = yt_tf.paragraphs[0]
        run = p.add_run()
        run.text = "🎥 Watch related video"
        run.hyperlink.address = yt_url
        run.font.size = Pt(16)
        run.font.color.rgb = font_color
        run.font.name = content_font
        p.alignment = PP_ALIGN.CENTER

    for i, slide in enumerate(slides[1:], 1):
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background_image(ppt_slide, prs, background_image_path)

        slide_width = prs.slide_width
        slide_height = prs.slide_height
        image_width = slide_width * 0.38
        text_width = slide_width * 0.58
        margin = Inches(0.4)
        
        # Clean the title for display and filename
        display_title = clean_title(slide['title'], language)
        # Check if this is the practice slide (last slide)
        is_practice_slide = "practice" in display_title.lower() or "question" in display_title.lower()

        # Title box (full width, but above text/image)
        title_box = ppt_slide.shapes.add_textbox(Inches(0.3), Inches(0.3), slide_width - Inches(0.6), Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.clear()
        title_tf.word_wrap = True
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_tf.margin_left = Inches(0.1)
        title_tf.margin_right = Inches(0.1)
        p = title_tf.paragraphs[0]
        p.text = display_title
        p.font.bold = True
        p.font.underline = True
        p.font.name = title_font
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
        
        # Dynamic font sizing for content slide titles
        title_len = len(display_title)
        if title_len > 100:
            p.font.size = Pt(20)
        elif title_len > 70:
            p.font.size = Pt(24)
        elif title_len > 40:
            p.font.size = Pt(28)
        else:
            p.font.size = Pt(32)

        # For practice slides, use full width content area (no image)
        if is_practice_slide:
            content_left = Inches(0.4)
            content_width = slide_width - Inches(0.8)
        else:
            if rtl:
                # RTL: image on left, text on right
                content_left = image_width + Inches(0.6)
            else:
                # LTR: text on left
                content_left = Inches(0.4)
            content_width = text_width - Inches(0.5)
        # Content text box (auto-size and dynamic font reduction)
        # Increase content area height to prevent overflow
        if language_key == 'urdu':
            content_height = slide_height - Inches(1.5)  # More space for Urdu content
        else:
            content_height = slide_height - Inches(1.8)  # More space for English content
        content_box = ppt_slide.shapes.add_textbox(content_left, Inches(1.1), content_width, content_height)
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        content_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        content_tf.margin_top = Inches(0.1)
        content_tf.margin_bottom = Inches(0.1)
        content_tf.margin_left = Inches(0.1)
        content_tf.margin_right = Inches(0.1)

        points = []
        for raw_line in slide['content']:
            points.extend(split_into_bullets(raw_line))

        # Limit bullets based on language - fewer for Urdu to prevent overflow
        if language_key == 'urdu':
            points = points[:4]  # Limit to 4 bullets for Urdu
        else:
            points = points[:6]  # Limit to 6 bullets for other languages
        point_count = len(points)
        # Dynamic font size reduction for overflow
        if language_key == 'urdu':
            # Larger font sizing for Urdu to ensure readability
            if is_practice_slide:
                font_size = Pt(28) if point_count <= 4 else Pt(24)
            else:
                if point_count <= 3:
                    font_size = Pt(32)
                elif point_count == 4:
                    font_size = Pt(28)
                elif point_count == 5:
                    font_size = Pt(24)
                elif point_count == 6:
                    font_size = Pt(20)
                else:
                    font_size = Pt(18)
            min_font_size = Pt(18)
        else:
            # Normal font sizing for other languages
            if is_practice_slide:
                font_size = Pt(20) if point_count <= 5 else Pt(18)
            else:
                if point_count <= 3:
                    font_size = Pt(24)
                elif point_count == 4:
                    font_size = Pt(20)
                elif point_count == 5:
                    font_size = Pt(18)
                elif point_count == 6:
                    font_size = Pt(16)
                else:
                    font_size = Pt(14)
        
        min_font_size = Pt(14) if language_key == 'urdu' else Pt(12)

        for idx, bullet in enumerate(points, 1):
            para = content_tf.add_paragraph()
            clean_bullet = bullet.lstrip('•- ').strip()
            if rtl and language_key == 'urdu':
                para.text = clean_bullet
            elif rtl:
                para.text = f"{clean_bullet} .{idx}"
            else:
                para.text = f"{idx}. {clean_bullet}"
            scaling = CONTENT_FONT_SCALING.get(content_font, 1.0)
            para.font.size = Pt(max(min_font_size.pt, font_size.pt * scaling))
            if para.font.size < min_font_size:
                para.font.size = min_font_size
            para.font.color.rgb = font_color
            para.font.name = content_font
            para.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
            # Increase line spacing for Urdu content
            if language_key == 'urdu':
                para.line_spacing = Pt(para.font.size.pt + 12)  # More spacing for Urdu
            else:
                para.line_spacing = Pt(para.font.size.pt + 6)
            para.level = 0  # This enables PowerPoint's own bullet
            if rtl:
                para.rtl = True

            # Simple text formatting without bold for content
            if not rtl:
                para.clear()
                bullet_text = f"{idx}. {clean_bullet}"
                run = para.add_run()
                run.text = bullet_text
                run.font.bold = False
                run.font.color.rgb = font_color
                run.font.name = content_font
                run.font.size = Pt(max(min_font_size.pt, font_size.pt * scaling))

        # Only add image if it's not a practice slide
        if not is_practice_slide:
            img_stream = fetch_unsplash_image(subject, display_title, language)
            if img_stream:
                try:
                    img_path = f"temp_img_{display_title.replace(' ', '_')}.jpg"
                    with open(img_path, "wb") as f:
                        f.write(img_stream.read())
                    # Vertically center the image in the right/left 38%
                    if rtl:
                        img_left = Inches(0.3)
                    else:
                        img_left = slide_width - image_width - Inches(0.3)
                    img_top = Inches(1.1) + (slide_height - Inches(2.0) - image_width) / 2
                    ppt_slide.shapes.add_picture(
                        img_path,
                        left=img_left,
                        top=img_top,
                        width=image_width,
                        height=image_width
                    )
                    os.remove(img_path)
                except Exception as e:
                    print(f"\u26a0\ufe0f Could not add image: {e}")

        footer = ppt_slide.shapes.add_textbox(margin, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5))
        footer_tf = footer.text_frame
        footer_text = {
            "english": "Generated by AI Slide Generator",
            "hindi": "एआई स्लाइड जनरेटर द्वारा निर्मित",
            "marathi": "एआई स्लाइड जनरेटर द्वारे तयार केले",
            "urdu": "سلائیڈ جنریٹر کے ذریعے تیار کردہ"
        }.get(language_key, "Generated by AI Slide Generator")
        footer_tf.text = footer_text
        footer_tf.paragraphs[0].font.size = Pt(12)
        footer_tf.paragraphs[0].font.color.rgb = font_color
        footer_tf.paragraphs[0].font.name = content_font
        footer_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = sanitize_filename(f"Class{class_level}_{subject}_{topic}_presentation_{timestamp}.pptx".replace(" ", ""))
    prs.save(filename)
    print(f"\u2705 PowerPoint saved as: {filename}")


def run_generator():
    print("\nAI Slide Generator")
    class_level = input("Grade level (e.g. 5): ")
    subject = input("Subject: ")
    topic = input("Topic: ")
    language = input("Language [default English]: ") or "English"
    try:
        num_slides = int(input("Number of slides [3-10, default 5]: ") or 5)
        if num_slides < 3 or num_slides > 10:
            raise ValueError
    except ValueError:
        print("\u274c Number of slides must be between 3 and 10.")
        return

    content = generate_ppt_content(class_level, subject, topic, language, num_slides)
    if isinstance(content, list):
        create_powerpoint(content, topic, class_level, subject, language)
    else:
        print(content)


if __name__ == "__main__":
    run_generator()