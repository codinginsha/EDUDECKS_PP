import streamlit as st
import google.generativeai as genai
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import os
import random
import re
import string
from googletrans import Translator
from googleapiclient.discovery import build
import datetime
import tempfile

# === API KEYS ===
GENAI_API_KEY = "AIzaSyDtX3Jdl8Kz_mDyrkdHBTWB3MNNwaXFEPo"
UNSPLASH_ACCESS_KEY_1 = "x7P3BqymG6BFDaFhnAiI_1ROJXsYt3U8xfpoLl5fEuM"
YOUTUBE_API_KEY = "AIzaSyD69g0uo0ldfuCrCKlEzCeMQao8QgrDGDY"

# === Configuration ===
genai.configure(api_key=GENAI_API_KEY)
translator = Translator()

# Simplified background colors for web deployment
BACKGROUND_COLORS = [
    "#FAFAFA",  # Soft White
    "#F5F5F5",  # Off-White
    "#EAF4FB",  # Light Pastel Blue
    "#DDEAF6",  # Gentle Sky Blue
    "#F9E6E6",  # Dusty Rose
    "#FCEEEF",  # Blush Pink
    "#FFF8E7",  # Pale Beige / Cream
    "#FDF6EC",  # Warm Cream
    "#E8F0E3",  # Muted Sage Green
    "#DDEBDD",  # Calm Green Tint
]

# Font configurations
TITLE_FONT_MAP = {
    "urdu": ["Arial", "Times New Roman"],
    "marathi": ["Arial", "Times New Roman"],
    "hindi": ["Arial", "Times New Roman"],
    "english": ["Arial", "Times New Roman", "Calibri", "Verdana"]
}

CONTENT_FONT_MAP = {
    "urdu": ["Arial", "Times New Roman"],
    "marathi": ["Arial", "Times New Roman"],
    "hindi": ["Arial", "Times New Roman"],
    "english": ["Arial", "Times New Roman", "Calibri", "Verdana"]
}

FONT_LIST = ["Arial", "Times New Roman", "Calibri", "Verdana"]

# Simple background images for web deployment
BACKGROUND_IMAGE_PATHS = [
    "https://images.unsplash.com/photo-1557804506-669a67965ba0?w=1200&h=800&fit=crop",
    "https://images.unsplash.com/photo-1557683316-973673baf926?w=1200&h=800&fit=crop",
    "https://images.unsplash.com/photo-1557682250-33bd709cbe85?w=1200&h=800&fit=crop",
    "https://images.unsplash.com/photo-1557683311-eac922347aa1?w=1200&h=800&fit=crop",
]

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def sanitize_filename(filename):
    valid_chars = f"-_.() {string.ascii_letters}{string.digits}"
    return ''.join(c for c in filename if c in valid_chars)

def translate_to_english(text, lang='auto'):
    try:
        translated = translator.translate(text, src=lang, dest='en')
        return translated.text
    except Exception as e:
        st.error(f"Translation error: {e}")
        return text

def fetch_unsplash_image(subject, query, language='English'):
    search_term = f"{subject} {query}"
    if language.lower() != 'english':
        search_term = translate_to_english(search_term, lang=language.lower())
    search_term = search_term.strip().replace(" ", "+")

    url = f"https://api.unsplash.com/photos/random?query={search_term}&orientation=landscape&client_id={UNSPLASH_ACCESS_KEY_1}"
    try:
        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            image_url = data['urls']['regular']
            st.success(f"Found image for: {search_term}")
            img_data = requests.get(image_url).content
            return BytesIO(img_data)
        else:
            st.warning(f"Unsplash error: {response.status_code}")
    except Exception as e:
        st.error(f"Image fetch error: {e}")
    return None

def search_youtube_video(query):
    search_query = query.replace(' ', '+')
    return f"https://www.youtube.com/results?search_query={search_query}"

def split_into_bullets(text):
    lines = text.split('\n')
    bullets = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("•") or line.startswith("-"):
            clean = re.sub(r"^[-•\s]+", "", line)
            bullets.append(f"• {clean}")
        else:
            parts = re.split(r'(?<=[.!?])\s+', line)
            for part in parts:
                part = part.strip()
                if part:
                    bullets.append(f"• {part}")
    return bullets

def clean_title(title, language=None):
    title = re.sub(r'[\\/*?<>|:"\n]', '', title)
    title = title.replace('<u>', '').replace('</u>', '').strip()
    
    if language and language.lower() == 'urdu':
        title = re.sub(r'^[uU]\s*', '', title)
        title = re.sub(r'\s*[uU]$', '', title)
        title = re.sub(r'^[A-Za-z\s\-\_\,\.\!\?\(\)\[\]]+', '', title)
        title = re.sub(r'[A-Za-z\s\-\_\,\.\!\?\(\)\[\]]+$', '', title)
        title = re.sub(r'[A-Za-z]+', '', title)
        title = title.strip()
    else:
        if title.lower().startswith('u') and len(title) > 1 and title[1].isupper():
            title = title[1:]
        if title.endswith('u') and len(title) > 2 and not title[-2].isspace():
            title = title[:-1]
        if title.endswith(' u'):
            title = title[:-2]
    return title.strip()

def generate_ppt_content(class_level, subject, topic, language='English', num_slides=5):
    model = genai.GenerativeModel('gemini-2.0-flash')
    is_language_subject = subject.strip().lower() in ["english", "hindi", "urdu", "marathi"]
    is_poem_or_lesson = any(x in topic.strip().lower() for x in ["poem", "lesson"])
    
    prompt = f"""Create a {num_slides}-slide PowerPoint presentation for a {class_level}th grade {subject} class on the topic "{topic}".
Use only formal, culturally accurate, and age-appropriate {language} language.
Do not mix English with {language}. No Roman script or SMS-style writing.
Each slide should include:
- A meaningful, localized title (bold and underlined using *<u>Title</u>* format).
- 5-6 clear bullet points using (• or -), with examples or facts.
- Do not write paragraphs, just concise bullet points.
- Avoid any unrelated historical figures or general knowledge.

If the subject or topic is related to math, mathematics, numericals, or formulas, include relevant formulas and worked examples on the slides. Use proper Unicode math formatting (e.g., x², fractions, π, √, etc.) and avoid ASCII-only equations. Show at least one worked example for each formula if possible.
"""
    
    if is_language_subject and is_poem_or_lesson:
        prompt += f"\nAfter the content slides, add one slide titled 'Difficult Words and Meanings' (in {language}) with 5-8 hard words from the poem or lesson and their meanings as bullet points. Format: • word: meaning."
    
    prompt += f"\nAfter the {{num_slides}} content slides, add one practice slide with 5 questions."
    prompt += f"\n\nRespond exactly in this format:\n\nSLIDE 1: [Title]\n[• Bullet or - Bullet or sentence]\n\nSLIDE 2: [Title]\n[• Bullet]\n[...]\n\nSLIDE {{num_slides + 1}}: Practice Questions\n[• Question 1]\n[• Question 2]\n[• Question 3]\n[• Question 4]\n[• Question 5]"
    
    try:
        with st.spinner("Generating presentation content..."):
            response = model.generate_content(prompt)
            return format_ppt_output(response.text)
    except Exception as e:
        st.error(f"Error generating content: {str(e)}")
        return None

def format_ppt_output(raw_text):
    slides = []
    current_slide = {}
    for line in raw_text.split('\n'):
        if line.strip().startswith('SLIDE'):
            if current_slide:
                slides.append(current_slide)
            parts = line.split(':', 1)
            current_slide = {'title': clean_title(parts[1].strip()), 'content': []}
        elif line.strip():
            current_slide['content'].append(line.strip())
    if current_slide:
        slides.append(current_slide)
    return slides

def create_powerpoint(slides, topic, class_level, subject, language="English"):
    prs = Presentation()
    language_key = language.lower().strip()
    rtl = True if language_key == 'urdu' else False

    # Pick fonts for this presentation
    title_font_choices = TITLE_FONT_MAP.get(language_key, FONT_LIST)
    content_font_choices = CONTENT_FONT_MAP.get(language_key, FONT_LIST)
    title_font = random.choice(title_font_choices)
    content_font = random.choice(content_font_choices)
    
    # Use a simple background color instead of image for web deployment
    background_color = random.choice(BACKGROUND_COLORS)
    font_color = RGBColor(0, 0, 0) if hex_to_rgb(background_color)[0] > 128 else RGBColor(255, 255, 255)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(*hex_to_rgb(background_color))

    first_title = clean_title(slides[0]['title'], language)
    slide_width = float(prs.slide_width) / 914400 if prs.slide_width else 10.0
    slide_height = float(prs.slide_height) / 914400 if prs.slide_height else 7.5

    box_height = 3.5
    box_top = (slide_height - box_height) / 2.0
    title_box = title_slide.shapes.add_textbox(
        Inches(1.0),
        Inches(box_top),
        Inches(slide_width - 2.0),
        Inches(box_height)
    )
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    p = tf.add_paragraph()
    p.text = first_title
    p.font.bold = True
    p.font.name = title_font
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.CENTER
    title_len = len(first_title)
    if title_len > 120:
        p.font.size = Pt(36)
    elif title_len > 80:
        p.font.size = Pt(44)
    elif title_len > 50:
        p.font.size = Pt(50)
    else:
        p.font.size = Pt(54)

    if slides[0]['content']:
        subtitle = slides[0]['content'][0]
        sp = tf.add_paragraph()
        sp.text = subtitle[:250]
        sp.font.size = Pt(32) if len(subtitle) < 80 else Pt(24)
        sp.font.color.rgb = font_color
        sp.font.name = content_font
        sp.alignment = PP_ALIGN.CENTER
        sp.space_before = Pt(24)

    # Content slides
    for i, slide in enumerate(slides[1:], 1):
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[6])
        ppt_slide.background.fill.solid()
        ppt_slide.background.fill.fore_color.rgb = RGBColor(*hex_to_rgb(background_color))

        slide_width = prs.slide_width
        slide_height = prs.slide_height
        text_width = slide_width * 0.9
        margin = Inches(0.4)
        
        display_title = clean_title(slide['title'], language)
        is_practice_slide = "practice" in display_title.lower() or "question" in display_title.lower()

        # Title box
        title_box = ppt_slide.shapes.add_textbox(Inches(0.3), Inches(0.3), slide_width - Inches(0.6), Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.clear()
        title_tf.word_wrap = True
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_tf.margin_left = Inches(0.1)
        title_tf.margin_right = Inches(0.1)
        p = title_tf.paragraphs[0]
        p.text = display_title
        p.font.bold = True
        p.font.underline = True
        p.font.name = title_font
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
        
        title_len = len(display_title)
        if title_len > 100:
            p.font.size = Pt(20)
        elif title_len > 70:
            p.font.size = Pt(24)
        elif title_len > 40:
            p.font.size = Pt(28)
        else:
            p.font.size = Pt(32)

        # Content area
        content_left = Inches(0.4)
        content_width = text_width - Inches(0.5)
        content_height = slide_height - Inches(1.8)
        content_box = ppt_slide.shapes.add_textbox(content_left, Inches(1.1), content_width, content_height)
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        content_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        content_tf.margin_top = Inches(0.1)
        content_tf.margin_bottom = Inches(0.1)
        content_tf.margin_left = Inches(0.1)
        content_tf.margin_right = Inches(0.1)

        points = []
        for raw_line in slide['content']:
            points.extend(split_into_bullets(raw_line))

        points = points[:6] if language_key != 'urdu' else points[:4]
        point_count = len(points)
        
        if language_key == 'urdu':
            font_size = Pt(28) if point_count <= 4 else Pt(24)
            min_font_size = Pt(18)
        else:
            font_size = Pt(20) if point_count <= 5 else Pt(18)
            min_font_size = Pt(14)

        for idx, bullet in enumerate(points, 1):
            para = content_tf.add_paragraph()
            clean_bullet = bullet.lstrip('•- ').strip()
            if rtl and language_key == 'urdu':
                para.text = clean_bullet
            elif rtl:
                para.text = f"{clean_bullet} .{idx}"
            else:
                para.text = f"{idx}. {clean_bullet}"
            
            para.font.size = Pt(max(min_font_size.pt, font_size.pt))
            para.font.color.rgb = font_color
            para.font.name = content_font
            para.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
            para.line_spacing = Pt(para.font.size.pt + 6)
            para.level = 0

        # Footer
        footer = ppt_slide.shapes.add_textbox(margin, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5))
        footer_tf = footer.text_frame
        footer_text = {
            "english": "Generated by EduDECKS AI",
            "hindi": "एडुडेक्स एआई द्वारा निर्मित",
            "marathi": "एडुडेक्स एआई द्वारे तयार केले",
            "urdu": "ایڈوڈیکس اے آئی کے ذریعے تیار کردہ"
        }.get(language_key, "Generated by EduDECKS AI")
        footer_tf.text = footer_text
        footer_tf.paragraphs[0].font.size = Pt(12)
        footer_tf.paragraphs[0].font.color.rgb = font_color
        footer_tf.paragraphs[0].font.name = content_font
        footer_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    # Save to temporary file
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = sanitize_filename(f"Class{class_level}_{subject}_{topic}_presentation_{timestamp}.pptx".replace(" ", ""))
    
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        prs.save(tmp_file.name)
        with open(tmp_file.name, 'rb') as f:
            pptx_data = f.read()
        os.unlink(tmp_file.name)
    
    return pptx_data, filename

# Streamlit App
def main():
    st.set_page_config(
        page_title="EduDECKS - AI PowerPoint Generator",
        page_icon="🎓",
        layout="wide"
    )
    
    st.title("🎓 EduDECKS - AI PowerPoint Generator")
    st.markdown("Create beautiful educational presentations with AI")
    
    # Sidebar for inputs
    with st.sidebar:
        st.header("📝 Presentation Settings")
        
        class_level = st.number_input("Grade Level", min_value=1, max_value=12, value=5, step=1)
        
        subject = st.text_input("Subject", placeholder="e.g., Science, Math, English")
        
        topic = st.text_input("Topic", placeholder="e.g., Photosynthesis, Algebra, Grammar")
        
        language = st.selectbox(
            "Language",
            ["English", "Hindi", "Urdu", "Marathi"],
            index=0
        )
        
        num_slides = st.slider("Number of Slides", min_value=3, max_value=10, value=5, step=1)
        
        st.markdown("---")
        st.markdown("**Features:**")
        st.markdown("• AI-generated content")
        st.markdown("• Multiple languages")
        st.markdown("• Beautiful templates")
        st.markdown("• Practice questions")
        st.markdown("• Download ready")
    
    # Main content area
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("📋 Presentation Details")
        
        if st.button("🚀 Generate Presentation", type="primary", use_container_width=True):
            if not subject or not topic:
                st.error("Please fill in both Subject and Topic fields.")
            else:
                with st.spinner("Creating your presentation..."):
                    # Generate content
                    content = generate_ppt_content(class_level, subject, topic, language, num_slides)
                    
                    if content:
                        # Create PowerPoint
                        pptx_data, filename = create_powerpoint(content, topic, class_level, subject, language)
                        
                        # Provide download button
                        st.success("✅ Presentation generated successfully!")
                        st.download_button(
                            label="📥 Download PowerPoint",
                            data=pptx_data,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
                        
                        # Show preview
                        st.subheader("📄 Content Preview")
                        for i, slide in enumerate(content[:3]):  # Show first 3 slides
                            with st.expander(f"Slide {i+1}: {slide['title']}"):
                                st.write("**Content:**")
                                for bullet in slide['content'][:3]:  # Show first 3 bullets
                                    st.write(f"• {bullet}")
                    else:
                        st.error("Failed to generate presentation content. Please try again.")
    
    with col2:
        st.subheader("🎨 Preview")
        st.markdown("""
        **Sample Presentation:**
        
        📊 **Title Slide**
        - Beautiful background
        - Professional fonts
        - Clean layout
        
        📝 **Content Slides**
        - AI-generated content
        - Bullet points
        - Relevant images
        
        ❓ **Practice Questions**
        - Interactive exercises
        - Learning reinforcement
        """)
        
        st.markdown("---")
        st.markdown("**Supported Languages:**")
        st.markdown("🇺🇸 English")
        st.markdown("🇮🇳 Hindi")
        st.markdown("🇵🇰 Urdu")
        st.markdown("🇮🇳 Marathi")

if __name__ == "__main__":
    main() 